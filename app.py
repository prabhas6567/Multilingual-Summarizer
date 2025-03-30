import streamlit as st  # type: ignore
from PyPDF2 import PdfReader  # type: ignore
import pandas as pd  # type: ignore
from pptx import Presentation  # type: ignore
from docx import Document  # type: ignore
import requests  # type: ignore
from bs4 import BeautifulSoup  # type: ignore
import os
import base64
from langchain.text_splitter import RecursiveCharacterTextSplitter  # type: ignore
from langchain_google_genai import GoogleGenerativeAIEmbeddings  # type: ignore
import google.generativeai as genai  # type: ignore
from langchain_community.vectorstores import FAISS  # type: ignore
from langchain_google_genai import ChatGoogleGenerativeAI  # type: ignore
from langchain.prompts import PromptTemplate  # type: ignore
from langchain.chains import RetrievalQA  # type: ignore
from dotenv import load_dotenv # type: ignore
from pytube import YouTube  # type: ignore
from youtube_transcript_api import YouTubeTranscriptApi  # type: ignore
import time  # type: ignore
from yt_dlp import YoutubeDL  # type: ignore
from PIL import Image  # type: ignore
import pytesseract # type: ignore
import easyocr  # type: ignore
from langdetect import detect  # type: ignore
from google.cloud import translate_v2 as translate
from transformers import pipeline # type: ignore

# Set the KMP_DUPLICATE_LIB_OK environment variable
os.environ['KMP_DUPLICATE_LIB_OK'] = 'TRUE'

# Load environment variables
load_dotenv('api.env')

# Configure Google GenAI
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    st.error("Google API Key is not configured! Please check your .env file.")
    st.stop()
genai.configure(api_key=api_key)

# List of 40 languages
languages = [
    "English","Afrikaans", "Albanian", "Amharic", "Arabic", "Armenian", "Azerbaijani", 
    "Bengali", "Bosnian", "Bulgarian", "Burmese","Catalan", "Cebuano", 
    "Chinese (Simplified)", "Chinese (Traditional)", "Croatian", "Czech",
    "Danish", "Dutch", "Estonian", "Filipino", "Finnish", 
    "French", "Galician", "Georgian", "German", "Greek", "Gujarati", 
    "Haitian Creole", "Hausa", "Hebrew", "Hindi", "Hungarian", "Icelandic", 
    "Indonesian", "Irish", "Italian", "Japanese", "Javanese", "Kannada", 
    "Kazakh", "Khmer", "Korean", "Kurdish", "Kyrgyz", "Lao", "Latin", "Latvian",
    "Lithuanian",  "Macedonian", "Malay", "Malayalam", "Maltese", "Marathi",
    "Mongolian", "Nepali", "Norwegian", "Persian", "Polish", "Portuguese", "Punjabi", 
    "Romanian", "Russian", "sanskrit","Serbian", "Sindhi", "Sinhala", "Slovak", "Slovenian",
    "Somali", "Spanish", "Swahili", "Swedish", "Tagalog", "Tajik", "Tamil",
    "Telugu", "Thai", "Turkish", "Ukrainian", "Urdu", "Uzbek", "Vietnamese", 
    "Welsh", "Yiddish", "Zulu",
]

# Function to get transcript from YouTube video
def get_transcript(video_id):
    transcript = YouTubeTranscriptApi.get_transcript(video_id)
    return " ".join([entry['text'] for entry in transcript])

# Function to detect language of the text
def detect_language(text):
    return detect(text)

# Function to translate text to the target language
def translate_text(text, target_language):
    translation = translate.Client().translate(text, target_language=target_language)
    return translation['translatedText']

# Function to summarize text using a pre-trained model
def summarize_text(text):
    summarizer = pipeline("summarization")
    return summarizer(text, max_length=130, min_length=30, do_sample=False)[0]['summary_text']

# Main function to summarize YouTube video
def summarize_youtube_video(video_id, target_language='en'):
    transcript = get_transcript(video_id)
    original_language = detect_language(transcript)

    if original_language != target_language:
        transcript = translate_text(transcript, target_language)

    summary = summarize_text(transcript)
    return summary

# Function to get file size in KB/MB
def get_file_size(file):
    file.seek(0, os.SEEK_END)  # Move to the end of the file
    size_bytes = file.tell()   # Get the file size in bytes
    file.seek(0)  # Reset the file pointer to the beginning
    size_kb = size_bytes / 1024
    if size_kb > 1024:
        return f"{size_kb / 1024:.2f} MB"
    return f"{size_kb:.2f} KB"

# Extract text from uploaded PDFs
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        num_pages = len(pdf_reader.pages)
        file_size = get_file_size(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

# Extract text from Excel files
def get_excel_text(excel_docs):
    text = ""
    for excel_file in excel_docs:
        # Use ExcelFile to get sheet names
        excel = pd.ExcelFile(excel_file)
        file_size = get_file_size(excel_file)
        st.write(f"**File Name:** {excel_file.name}")
        st.write(f"**File Size:** {file_size}")
        st.write(f"**Number of Sheets:** {len(excel.sheet_names)}")
        
        # Read each sheet and append its content to the text
        for sheet_name in excel.sheet_names:
            df = excel.parse(sheet_name)
            text += f"Sheet: {sheet_name}\n"
            text += df.to_string(index=False) + "\n\n"
    return text

# Extract text from PPT files
def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        presentation = Presentation(ppt)
        file_size = get_file_size(ppt)
        st.write(f"**File Name:** {ppt.name}")
        st.write(f"**File Size:** {file_size}")
        st.write(f"**Number of Slides:** {len(presentation.slides)}")
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

# Extract text from DOC files
def get_doc_text(doc_docs):
    text = ""
    for doc in doc_docs:
        document = Document(doc)
        file_size = get_file_size(doc)
        st.write(f"**File Name:** {doc.name}")
        st.write(f"**File Size:** {file_size}")
        st.write(f"**Number of Paragraphs:** {len(document.paragraphs)}")
        for para in document.paragraphs:
            text += para.text + "\n"
    return text

# Extract text from TXT files
def get_txt_text(txt_docs):
    text = ""
    for txt_file in txt_docs:
        file_size = get_file_size(txt_file)
        st.write(f"**File Name:** {txt_file.name}")
        st.write(f"**File Size:** {file_size}")
        content = txt_file.read().decode("utf-8")
        st.write(f"**Content Preview:** {content[:100]}...")  # Show a preview of the content
        text += content + "\n"
    return text

# Extract text from a website
def get_website_text(url):
    try:
        response = requests.get(url)
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Extract title
        title = soup.title.string if soup.title else "No Title"
        
        # Extract main content (e.g., paragraphs, headings)
        text = ""
        for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
            text += element.get_text() + "\n"
        
        return {
            "title": title,
            "content": text.strip()  # Remove any leading/trailing whitespace
        }
    except Exception as e:
        return {
            "title": "Error",
            "content": f"Could not retrieve website content: {e}"
        }

# Extract details and transcript from a YouTube video
def get_youtube_details_and_transcript(video_url):
    ydl_opts = {
        'quiet': True,
        'extract_flat': True,
    }
    with YoutubeDL(ydl_opts) as ydl:
        try:
            info = ydl.extract_info(video_url, download=False)
            title = info.get('title', 'Unknown Title')
            channel = info.get('uploader', 'Unknown Channel')
            thumbnail_url = info.get('thumbnail', '')
            video_id = info.get('id', '')
            views = info.get('view_count', 'Unknown Views')
            likes = info.get('like_count', 'Unknown Likes')
            dislikes = info.get('dislike_count', 'Unknown Dislikes')
            comments = info.get('comment_count', 'Unknown Comments')
            duration = info.get('duration', 'Unknown Duration')
            description = info.get('description', 'Unknown Description')

            # Get transcript
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['te'])
                transcript_text = ' '.join([entry['text'] for entry in transcript])
            except Exception as e:
                transcript_text = f"Could not retrieve transcript: {e}"

            return {
                "title": title,
                "channel": channel,
                "thumbnail_url": thumbnail_url,
                "video_id": video_id,
                "views": views,
                "likes": likes,
                "dislikes": dislikes,
                "comments": comments,
                "duration": duration,
                "description": description,
                "transcript": transcript_text
            }
        except Exception as e:
            st.error(f"Error retrieving YouTube video details: {e}")
            return {
                "title": "Unknown Title",
                "channel": "Unknown Channel",
                "thumbnail_url": "",
                "video_id": "",
                "views": "Unknown Views",
                "likes": "Unknown Likes",
                "dislikes": "Unknown Dislikes",
                "comments": "Unknown Comments",
                "duration": "Unknown Duration",
                "description": "Unknown Description",
                "transcript": "Could not retrieve transcript due to an error."
            }

# Split the extracted text into smaller chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
    chunks = text_splitter.split_text(text)
    return chunks

# Create and save the vector store
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# # Create the conversational chain for Q&A
# def get_conversational_chain(retriever, language, summary_length):
#     length_mapping = {
#         "Short": "Provide a brief summary.",
#         "Medium": "Provide a detailed summary.",
#         "Large": "Provide an extensive summary."
#     }
#     prompt_template = f"""Answer the question as detailed as possible from the provided context in {language}.  
# Make sure to provide all the details.  
# If the answer is not in the provided context, just say, "Answer is not available in the context."  
# Do not provide a wrong answer.  

# Context:  
# {{context}}  

# Question:  
# {{question}}  

# Answer:  
# {length_mapping.get(summary_length, "Provide a detailed summary.")}
# """
#     prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
#     model = ChatGoogleGenerativeAI(model="gemini-1.5-pro", temperature=0.3)

#     chain = RetrievalQA.from_chain_type(
#         llm=model,
#         retriever=retriever,
#         chain_type="stuff", 
#         chain_type_kwargs={"prompt": prompt},
#     )
#     return chain

# Create the conversational chain for Q&A
def get_conversational_chain(retriever, language, summary_length):
    length_mapping = {
        "Short": "Provide a brief summary.",
        "Medium": "Provide a detailed summary.",
        "Large": "Provide an extensive summary."
    }
    prompt_template = f"""Answer the question as detailed as possible from the provided context in {language}.  
    Make sure to provide all the details.  
    If the answer is not in the provided context, just say, "Answer is not available in the context."  
    Do not provide a wrong answer.  

    Context:  
    {{context}}  

    Question:  
    {{question}}  

    Answer:  
    {length_mapping.get(summary_length, "Provide a detailed summary.")}
    """
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    model = ChatGoogleGenerativeAI(model="gemini-1.5-pro", temperature=0.3)

    chain = RetrievalQA.from_chain_type(
        llm=model,
        retriever=retriever,
        chain_type="stuff", 
        chain_type_kwargs={"prompt": prompt},
    )
    return chain
# Maintain chat history
chat_history = []
def user_input(user_question, sample_text, language, summary_length):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    try:
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        retriever = new_db.as_retriever()
        chain = get_conversational_chain(retriever, language, summary_length)

        for attempt in range(5):  # Retry up to 5 times
            try:
                response = chain.invoke({"query": user_question})
                chat_history.append({"assistant": response['result']})
                break  # Exit loop if successful
            except Exception as e:
                if "ResourceExhausted" in str(e):
                    wait_time = 2 ** attempt  # Exponential backoff
                    st.warning(f"Quota exceeded, retrying in {wait_time} seconds...")
                    time.sleep(wait_time)
                else:
                    st.error(f"An error occurred: {e}")
                    break
    except Exception as e:
        st.error(f"An error occurred: {e}")

# Display chat history
def display_chat():
    for chat in chat_history:
        st.write(f"🤖 **Aura:** {chat['assistant']}")
        st.markdown("---")  # Separator for chat messages

# Main application logic
def main():
    st.set_page_config(
        page_title="A MultiLingual Summarizer", 
        page_icon="🤖",
        layout="wide")
    st.title("A MultiLingual Summarizer for Enhanced Text Comprehension")
    st.header("Welcome to Aura  - Your Summarizer Agent")

    # Define file-specific prompts
    file_specific_prompts = {
        "PDF": {
            "Summarize the given document": "Please summarize the following PDF document",
            "What are the key points?": "What are the key points of the PDF?",
            " Explain this concept": "Can you explain this concept from the PDF in simple terms?",
            "What's the conclusion?": "What is the conclusion or final thought of the PDF?",
            "What's the main idea?": "What is the main idea or topic of the PDF?",
            "What's the purpose?": "What is the purpose of the PDF?",
        },
        
        "Excel": {
            "Summarize this spreadsheet": "Please summarize the following Excel spreadsheet.",
            "What's the number of sheets?": "What is the number of sheets in the Excel file?",
            " Mention all the rows names given in the spreadsheet?": "Please Mention all the rows given in the spreadsheet",
            "What's the data type?": "What is the data type of the Excel file?",
        },
        
        "PPT": {
            "Summarize this presentation": "Please summarize the following PowerPoint presentation.",
            "What are the key points?": "What are the key points of the presentation?",
            "Summarize and Explain the  slide-1?": "Can you explain this slide in the presentation?",
             "Summarize and Explain the slide-2?": "Can you explain this slide in the presentation?",
            "Explain and summarize all the Slides  in brief?" : "Please Explain and summarize all the Slidesin brief?",
            "Describe all the contents given in PPT in pointwise?":"PLease Describe all the contents given in PPT in pointwise?",
        },
        
        "DOC": {
            "Summarize this document": "Please summarize the following Word document.",
            "What are the key points?": "What are the key points of the document?",
            "Explain this section": "Can you explain this section in the document?",
        },
        
        "TXT": {
            "Summarize this text": "Please summarize the following text file.",
            "What are the key points?": "What are the key points of the text?",
            "Explain this paragraph": "Can you explain this paragraph in the text?",
        },
        
        "YouTube": {
            "Summarize this video": "Please summarize the following YouTube video.",
            "What's the purpose?": "What is the purpose of the video?",
            "What are the key points?": "What are the key points of the video?",
            "Explain this part": "Can you explain this part of the video?",
            "Can you prepare notes based on the video?": "Prepare the notes based on the video",
        },
        "Website": {
            "Summarize this website": "Please summarize the following website content.",
            "What are the key points?": "What are the key points of the website?",
            "Explain this section": "Can you explain this section of the website?",
        },
        "Sample Text": {
            "Summarize this text": "Please summarize the following text.",
            "What are the key points?": "What are the key points of the text?",
            "Can you Summarize the given content and Detect the language given":"Please summarize the content given and also detect the language given of the text?",
            "Explain this paragraph": "Can you explain this paragraph in the text?",
        },
    }
     
   
    with st.sidebar:
        st.title("Menu:")
        file_type = st.selectbox("Select File Type", ["PDF", "Excel", "PPT", "DOC", "TXT",  "YouTube", "Website", "Sample Text"])
        
        # Language selection
        selected_language = st.selectbox("Select Language for Summarization", languages)

        # Summary length selection
        summary_length = st.selectbox("Select Summary Length", ["Short", "Medium", "Large"])

        raw_text = ""
 
        # PDF Preview and Details
        if file_type == "PDF":
           pdf_docs = st.file_uploader("Upload Your PDFs", type=["pdf"], accept_multiple_files=True)
           if pdf_docs:
              st.markdown("<div style='background-color: rgba(255, 255, 255, 0.8); padding: 10px; border-radius: 5px;'>", unsafe_allow_html=True)
              st.write("**PDF Details and Preview:**")
              for pdf in pdf_docs:
                  st.write(f"**File Name:** {pdf.name}")
                  st.write(f"**File Size:** {get_file_size(pdf)}")
                  pdf_reader = PdfReader(pdf)
                  st.write(f"**Number of Pages:** {len(pdf_reader.pages)}")
                  st.write("**Preview:**")
                  for page in pdf_reader.pages:
                      st.write(page.extract_text())
              st.markdown("</div>", unsafe_allow_html=True)

        # Excel Preview and Details
        elif file_type == "Excel":
             excel_docs = st.file_uploader("Upload Your Excel Files", type=["xlsx"], accept_multiple_files=True)
             if excel_docs:
                st.markdown("<div style='background-color: rgba(255, 255, 255, 0.8); padding: 10px; border-radius: 5px;'>", unsafe_allow_html=True)
                st.write("**Excel Details and Preview:**")
                for excel_file in excel_docs:
                    st.write(f"**File Name:** {excel_file.name}")
                    st.write(f"**File Size:** {get_file_size(excel_file)}")
                    excel = pd.ExcelFile(excel_file)
                    st.write(f"**Number of Sheets:** {len(excel.sheet_names)}")
                    st.write("**Preview:**")
                    for sheet_name in excel.sheet_names:
                        df = excel.parse(sheet_name)
                        st.write(f"**Sheet: {sheet_name}**")
                        st.dataframe(df, height=500)  # Display the full sheet with a scrollable area
                st.markdown("</div>", unsafe_allow_html=True)

       # PPT Preview and Details
        elif file_type == "PPT":
            ppt_docs = st.file_uploader("Upload Your PPTs", type=["pptx"], accept_multiple_files=True)
            if ppt_docs:
              st.markdown("<div style='background-color: rgba(255, 255, 255, 0.8); padding: 10px; border-radius: 5px;'>", unsafe_allow_html=True)
              st.write("**PPT Details and Preview:**")
              for ppt in ppt_docs:
                 st.write(f"**File Name:** {ppt.name}")
                 st.write(f"**File Size:** {get_file_size(ppt)}")
                 presentation = Presentation(ppt)
                 st.write(f"**Number of Slides:** {len(presentation.slides)}")
                 st.write("**Preview:**")
                 for slide in presentation.slides:
                     for shape in slide.shapes:
                         if hasattr(shape, "text"):
                            st.write(shape.text)
              st.markdown("</div>", unsafe_allow_html=True)   

       # DOC Preview and Details
        elif file_type == "DOC":
             doc_docs = st.file_uploader("Upload Your DOCs", type=["docx"], accept_multiple_files=True)
             if doc_docs:
                st.markdown("<div style='background-color: rgba(255, 255, 255, 0.8); padding: 10px; border-radius: 5px;'>", unsafe_allow_html=True)
                st.write("**DOC Details and Preview:**")
                for doc in doc_docs:
                    st.write(f"**File Name:** {doc.name}")
                    st.write(f"**File Size:** {get_file_size(doc)}")
                    document = Document(doc)
                    st.write(f"**Number of Paragraphs:** {len(document.paragraphs)}")
                    st.write("**Preview:**")
                    for para in document.paragraphs:
                        st.write(para.text)
                st.markdown("</div>", unsafe_allow_html=True)

# TXT Preview and Details
        elif file_type == "TXT":
            txt_docs = st.file_uploader("Upload Your TXT Files", type=["txt"], accept_multiple_files=True)
            if txt_docs:
               st.markdown("<div style='background-color: rgba(255, 255, 255, 0.8); padding: 10px; border-radius: 5px;'>", unsafe_allow_html=True)
               st.write("**TXT Details and Preview:**")
               for txt_file in txt_docs:
                   st.write(f"**File Name:** {txt_file.name}")
                   st.write(f"**File Size:** {get_file_size(txt_file)}")
                   content = txt_file.read().decode("utf-8")
                   st.write("**Preview:**")
                   st.write(content[:500] + "...")  # Show a preview of the first 500 characters
               st.markdown("</div>", unsafe_allow_html=True)

        elif file_type == "YouTube":
            youtube_link = st.text_input("Enter YouTube Video Link")
            if youtube_link:
                youtube_details = get_youtube_details_and_transcript(youtube_link)
                raw_text += youtube_details['transcript']
               
               

        elif file_type == "Website":
            website_link = st.text_input("Enter Website URL")

        elif file_type == "Sample Text":
            sample_text_input = st.text_area("Enter your sample text here:", height=150)

        # Move the Submit & Process button under the file upload section
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                try:
                    if 'pdf_docs' in locals() and pdf_docs:
                        raw_text += get_pdf_text(pdf_docs)
                    if 'excel_docs' in locals() and excel_docs:
                        raw_text += get_excel_text(excel_docs)
                    if 'ppt_docs' in locals() and ppt_docs:
                        raw_text += get_ppt_text(ppt_docs)
                    if 'doc_docs' in locals() and doc_docs:
                        raw_text += get_doc_text(doc_docs)
                    if 'txt_docs' in locals() and txt_docs:
                        raw_text += get_txt_text(txt_docs)
                    if 'youtube_link' in locals() and youtube_link:
                        if 'youtube_link' in locals() and youtube_link:
                            youtube_details = get_youtube_details_and_transcript(youtube_link)
                            st.image(youtube_details['thumbnail_url'], caption=f"{youtube_details['title']} by {youtube_details['channel']}")
                            st.write(f"**Video Title:** {youtube_details['title']}")
                            st.write(f"**Channel:** {youtube_details['channel']}")
                            st.write(f"**Video ID:** {youtube_details['video_id']}")
                            st.write(f"**Views:** {youtube_details['views']}")
                            st.write(f"**Likes:** {youtube_details['likes']}")
                            st.write(f"**Dislikes:** {youtube_details['dislikes']}")
                            st.write(f"**Comments:** {youtube_details['comments']}")
                            st.write(f"**Duration:** {youtube_details['duration']} seconds")
                            st.write(f"**Description:** {youtube_details['description']}")
                            raw_text += youtube_details['transcript']
                    if 'website_link' in locals() and website_link:
                        website_data = get_website_text(website_link)
                        raw_text += website_data['content']
                        st.write(f"**Website Title:** {website_data['title']}")
                        st.write(f"**Website Content:** {website_data['content']}")
                    if 'sample_text_input' in locals() and sample_text_input:
                        raw_text += sample_text_input
                    
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks)
                    st.success("Processing complete! You can now ask questions.")
                except Exception as e:
                    st.error(f"An error occurred while processing the files: {e}")

    # Update example prompts based on selected file type
    example_prompts = file_specific_prompts.get(file_type, {})
    selected_prompt = st.selectbox("Select a prompt:", list(example_prompts.keys()))
    user_question = st.text_input("What would you like to know?", placeholder=example_prompts[selected_prompt])  # Populate with selected prompt

    if st.button("Submit"):
        sample_text = sample_text_input if 'sample_text_input' in locals() and file_type == "Sample Text" else ""
        user_input(user_question, sample_text, selected_language, summary_length)
    display_chat()  # Display chat history

if __name__ == "__main__":
    main()
